import json
import logging
import os
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Union
from uuid import uuid4

from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


class OutputGenerator:
    """
    Generates narration artifacts (json/txt/docx/pptx) to a writable output directory.

    Defaults to /tmp/temp_outputs (container-friendly). Override with OUTPUT_DIR env var
    or by passing output_dir explicitly.
    """

    def __init__(self, output_dir: Union[str, Path, None] = None):
        output_dir = output_dir or os.getenv("OUTPUT_DIR", "/tmp/temp_outputs")
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    # ---------- helpers ----------

    def _safe_base(self, base_name: str) -> str:
        return "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in base_name).strip("_") or "output"

    def _run_suffix(self) -> str:
        # Unique across concurrent requests (timestamp + random)
        ts = datetime.utcnow().strftime("%Y%m%dT%H%M%S.%fZ")
        return f"{ts}_{uuid4().hex[:8]}"

    def _atomic_write_text(self, path: Path, content: str) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        fd, tmp = tempfile.mkstemp(prefix=path.stem + "_", suffix=path.suffix, dir=str(path.parent))
        tmp_path = Path(tmp)

        try:
            with os.fdopen(fd, "w", encoding="utf-8") as f:
                f.write(content)
                f.flush()
                os.fsync(f.fileno())
            tmp_path.replace(path)  # atomic on same filesystem
        finally:
            # In case replace failed
            try:
                if tmp_path.exists():
                    tmp_path.unlink()
            except OSError:
                pass

    def _validate_slides(self, result: Dict[str, Any]) -> List[Dict[str, Any]]:
        slides = result.get("slides", [])
        if not isinstance(slides, list):
            raise ValueError("result['slides'] must be a list")

        norm: List[Dict[str, Any]] = []
        for i, s in enumerate(slides):
            if not isinstance(s, dict):
                raise ValueError(f"slides[{i}] must be a dict")

            slide_number = s.get("slide_number")
            slide_number_int: int | None
            try:
                slide_number_int = int(slide_number)
            except Exception:
                slide_number_int = None

            norm.append(
                {
                    "slide_number": slide_number_int,
                    "narration_paragraph": str(s.get("narration_paragraph", "") or ""),
                    "speaker_notes": str(s.get("speaker_notes", "") or ""),
                }
            )
        return norm

    # ---------- generators ----------

    def generate_json(self, result: Dict[str, Any], base_name: str) -> Path:
        base = self._safe_base(base_name)
        filename = f"{base}_{self._run_suffix()}_narration.json"
        file_path = self.output_dir / filename
        try:
            payload = json.dumps(result, indent=2, ensure_ascii=False)
            self._atomic_write_text(file_path, payload)
            logger.info("Generated JSON output: %s", file_path)
            return file_path
        except Exception:
            logger.exception("Failed to generate JSON output")
            raise

    def generate_text(self, result: Dict[str, Any], base_name: str) -> Path:
        base = self._safe_base(base_name)
        filename = f"{base}_{self._run_suffix()}_narration.txt"
        file_path = self.output_dir / filename
        try:
            slides = self._validate_slides(result)

            lines: List[str] = [
                f"Narration Script for {base_name}",
                "=" * 50,
                "",
            ]

            for slide in slides:
                slide_num = slide["slide_number"] if slide["slide_number"] is not None else "?"
                lines.extend(
                    [
                        f"Slide {slide_num}",
                        "-" * 20,
                        "Narration:",
                        slide["narration_paragraph"],
                        "",
                        "Speaker Notes:",
                        slide["speaker_notes"],
                        "",
                        "=" * 50,
                        "",
                    ]
                )

            self._atomic_write_text(file_path, "\n".join(lines))
            logger.info("Generated Text output: %s", file_path)
            return file_path
        except Exception:
            logger.exception("Failed to generate Text output")
            raise

    def generate_word(self, result: Dict[str, Any], base_name: str) -> Path:
        base = self._safe_base(base_name)
        filename = f"{base}_{self._run_suffix()}_narration.docx"
        file_path = self.output_dir / filename
        try:
            slides = self._validate_slides(result)

            doc = Document()
            doc.add_heading(f"Narration Script for {base_name}", 0)

            for idx, slide in enumerate(slides):
                slide_num = slide["slide_number"] if slide["slide_number"] is not None else "?"
                doc.add_heading(f"Slide {slide_num}", level=1)

                doc.add_heading("Narration:", level=2)
                doc.add_paragraph(slide["narration_paragraph"])

                doc.add_heading("Speaker Notes:", level=2)
                doc.add_paragraph(slide["speaker_notes"])

                if idx != len(slides) - 1:
                    doc.add_page_break()

            # python-docx doesn't provide atomic save; keep output dir reliable instead.
            doc.save(file_path)
            logger.info("Generated Word output: %s", file_path)
            return file_path
        except Exception:
            logger.exception("Failed to generate Word output")
            raise

    def generate_pptx_with_notes(
        self,
        original_pptx: Union[str, Path],
        result: Dict[str, Any],
        base_name: str,
        mode: str = "replace",  # "replace" or "append"
    ) -> Path:
        """
        mode="replace": overwrite notes with narration only
        mode="append": append narration under a marker, preserving existing notes text
        """
        base = self._safe_base(base_name)
        filename = f"{base}_{self._run_suffix()}_with_narration.pptx"
        file_path = self.output_dir / filename
        try:
            original_pptx = Path(original_pptx)
            slides = self._validate_slides(result)

            narration_map: Dict[int, str] = {}
            for s in slides:
                if s["slide_number"] is not None:
                    narration_map[s["slide_number"]] = (s.get("narration_paragraph") or "").strip()

            prs = Presentation(str(original_pptx))

            marker = "--- Generated ---"

            for i, slide in enumerate(prs.slides, start=1):
                narration = narration_map.get(i)
                if not narration:
                    continue

                notes_slide = slide.notes_slide
                tf = notes_slide.notes_text_frame

                if mode == "append":
                    existing = (tf.text or "").strip()
                    tf.text = (
                        f"{existing}\n\n{marker}\n{narration}".strip()
                        if existing
                        else narration
                    )
                else:
                    tf.text = narration

            prs.save(str(file_path))
            logger.info("Generated PPTX output: %s", file_path)
            return file_path
        except Exception:
            logger.exception("Failed to generate PPTX output")
            raise

    # ---------- optional cleanup utility ----------

    def cleanup_old_files(self, older_than_seconds: int = 6 * 3600) -> int:
        """
        Best-effort cleanup for /tmp output dirs. Deletes files older than threshold.
        Returns number of deleted files.
        """
        deleted = 0
        cutoff = datetime.utcnow().timestamp() - older_than_seconds
        try:
            for p in self.output_dir.glob("*"):
                try:
                    if p.is_file() and p.stat().st_mtime < cutoff:
                        p.unlink()
                        deleted += 1
                except OSError:
                    continue
        except OSError:
            pass
        return deleted