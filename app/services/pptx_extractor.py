import os
import shutil
import subprocess
import re
import logging
from pathlib import Path
from typing import List, Tuple, Optional
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

class PPTXExtractor:
    def __init__(self, temp_dir: Path):
        self.temp_dir = temp_dir
        self.temp_files = []

    def _resolve_soffice_cmd(self) -> str:
        """Resolve LibreOffice executable robustly (soffice may not be on PATH)."""
        soffice_cmd = (
            shutil.which("soffice")
            or shutil.which("libreoffice")
            or shutil.which("soffice.bin")
        )

        if not soffice_cmd:
            candidates = [
                "/usr/bin/soffice",
                "/usr/local/bin/soffice",
                "/usr/lib/libreoffice/program/soffice",
                "/usr/lib64/libreoffice/program/soffice",
                "/opt/libreoffice/program/soffice",
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            ]
            for c in candidates:
                if os.path.exists(c) and os.access(c, os.X_OK):
                    soffice_cmd = c
                    break

        if not soffice_cmd:
            raise FileNotFoundError(
                "LibreOffice executable not found. "
                "Install LibreOffice in the runtime image or ensure 'soffice'/'libreoffice' is on PATH."
            )

        return soffice_cmd

    def _extract_slide_text(self, slide) -> str:
        """Extract text content from a slide."""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_parts.append(text)
        return "\n".join(text_parts)

    def _extract_speaker_notes(self, slide) -> str:
        """Extract speaker notes from a slide."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    return notes_text if notes_text else ""
            return ""
        except Exception:
            # If notes slide doesn't exist or can't be accessed, return empty string
            return ""

    def pptx_to_images(self, pptx_path: Path) -> Tuple[List[Path], List[str], List[str]]:
        """
        Convert PPTX slides to images via PDF intermediate and extract speaker notes and original text.
        Returns tuple of (image_paths, speaker_notes_list, original_text_list).
        """
        try:
            logger.info(f"Loading PowerPoint file: {pptx_path}")
            prs = Presentation(pptx_path)
            total_slides = len(prs.slides)
            logger.info(f"Found {total_slides} slides in presentation")

            image_paths = []
            speaker_notes_list = []
            original_text_list = []

            # Extract notes and text first
            for slide in prs.slides:
                speaker_notes_list.append(self._extract_speaker_notes(slide))
                original_text_list.append(self._extract_slide_text(slide))

            # Check dependencies first
            soffice_path = shutil.which("soffice")
            pdftoppm_path = shutil.which("pdftoppm")
            logger.info(f"External tools check: soffice={'FOUND' if soffice_path else 'MISSING'}, pdftoppm={'FOUND' if pdftoppm_path else 'MISSING'}")

            # Attempt Image Conversion
            try:
                logger.info("Converting PPTX to PDF using LibreOffice (soffice)...")
                
                # Create directories
                pdf_dir = self.temp_dir / "rendered_pdf"
                pdf_dir.mkdir(parents=True, exist_ok=True)
                self.temp_files.append(pdf_dir)
                
                png_dir = self.temp_dir / "rendered_png"
                png_dir.mkdir(parents=True, exist_ok=True)
                self.temp_files.append(png_dir)

                soffice_cmd = self._resolve_soffice_cmd()
                
                # Step 1: PPTX -> PDF
                pdf_out = pdf_dir / f"{pptx_path.stem}.pdf"
                self.temp_files.append(pdf_out)

                subprocess.run(
                    [
                        soffice_cmd,
                        "--headless",
                        "--nologo",
                        "--nolockcheck",
                        "--nodefault",
                        "--norestore",
                        "--convert-to",
                        "pdf:impress_pdf_Export",
                        "--outdir",
                        str(pdf_dir),
                        str(pptx_path),
                    ],
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    check=True,
                )

                # Find the generated PDF file
                pdf_candidates = list(pdf_dir.glob("*.pdf"))
                if not pdf_candidates:
                    raise Exception("LibreOffice did not produce a PDF file")
                pdf_out = pdf_candidates[0]
                self.temp_files.append(pdf_out)

                logger.info(f"PDF created: {pdf_out}")
                logger.info("Converting PDF to PNG using pdftoppm...")

                # Step 2: PDF -> PNG using pdftoppm
                pdftoppm_cmd = shutil.which("pdftoppm")
                if not pdftoppm_cmd:
                    raise FileNotFoundError(
                        "pdftoppm not found. Install 'poppler-utils' to enable PDF-to-PNG conversion."
                    )

                # Generate PNG files with prefix
                out_prefix = str(png_dir / "slide")
                result = subprocess.run(
                    [
                        pdftoppm_cmd,
                        "-png",
                        "-r", "150",  # Lower resolution
                        "-scale-to-x", "1280",  # Constrain width
                        "-scale-to-y", "-1",  # Maintain aspect ratio
                        str(pdf_out),
                        out_prefix,
                    ],
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    check=True,
                )

                # Get generated PNG files
                png_files = list(png_dir.glob("*.png"))
                
                def _slide_sort_key(p: Path):
                    m = re.search(r"(\d+)(?=\D*$)", p.stem)
                    return int(m.group(1)) if m else 10**9

                png_files = sorted(png_files, key=_slide_sort_key)
                
                # Take only the number of slides we expect
                if len(png_files) > total_slides:
                    png_files = png_files[:total_slides]
                
                image_paths = [p for p in png_files]
                for p in image_paths:
                    self.temp_files.append(p)

                logger.info(f"Successfully converted {len(image_paths)} slides to images")
                
                return image_paths, speaker_notes_list, original_text_list

            except subprocess.CalledProcessError as e:
                error_msg = f"Subprocess failed: {e.cmd}"
                if e.stderr:
                    error_msg += f"\nStderr: {e.stderr.decode()}"
                if e.stdout:
                    error_msg += f"\nStdout: {e.stdout.decode()}"
                logger.error(error_msg)
                # Return partial results (text/notes) even if images fail
                return [], speaker_notes_list, original_text_list

            except Exception as e:
                logger.error(f"Image generation failed (continuing without images): {str(e)}")
                import traceback
                logger.error(traceback.format_exc())
                return [], speaker_notes_list, original_text_list
                # We return empty image_paths, but still return text/notes
                # The processor must handle the case where image_paths is empty or mismatched

            # If image generation failed, we might have 0 images. 
            # If it succeeded but slide count mismatch, that's handled by processor.
            return image_paths, speaker_notes_list, original_text_list

        except Exception as e:
            logger.error(f"Failed to extract content from PPTX: {e}")
            raise Exception(f"Failed to extract content from PPTX: {str(e)}")