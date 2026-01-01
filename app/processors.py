import os
import tempfile
import json
import logging
import subprocess
import re
import shutil  # ✅ needed for shutil.which
from pathlib import Path
from typing import List, Dict, Optional, Tuple
import google.generativeai as genai
from pptx import Presentation
from PIL import Image
import io
import base64

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
)
logger = logging.getLogger(__name__)


class SlideProcessor:
    def __init__(self, api_key: str, model_name: str = "gemini-2.5-flash"):
        """Initialize the slide processor with Gemini API."""
        genai.configure(api_key=api_key)
        self.model = genai.GenerativeModel(model_name)
        self.temp_dir = None
        self.temp_files = []

    def _extract_response_text(self, response) -> str:
        """
        Safely extract text from Gemini API response.
        Handles both simple text responses and complex multi-part responses.
        """
        # Try simple text access first (may raise exception for complex responses)
        try:
            return response.text
        except (AttributeError, ValueError, Exception) as e:
            logger.debug(f"Simple text access failed: {e}, trying parts accessor")

        # Fallback to parts accessor
        try:
            if hasattr(response, "parts") and response.parts:
                text_parts = []
                for part in response.parts:
                    if hasattr(part, "text") and part.text:
                        text_parts.append(part.text)
                if text_parts:
                    logger.debug("Successfully extracted text from response.parts")
                    return "".join(text_parts)
        except Exception as e:
            logger.warning(f"Error extracting from parts: {e}")

        # Fallback to candidates accessor
        try:
            if hasattr(response, "candidates") and response.candidates:
                for candidate in response.candidates:
                    if hasattr(candidate, "content") and candidate.content:
                        if hasattr(candidate.content, "parts") and candidate.content.parts:
                            text_parts = []
                            for part in candidate.content.parts:
                                if hasattr(part, "text") and part.text:
                                    text_parts.append(part.text)
                            if text_parts:
                                logger.debug("Successfully extracted text from response.candidates")
                                return "".join(text_parts)
        except Exception as e:
            logger.warning(f"Error extracting from candidates: {e}")

        # Last resort: try to get text from result object directly
        try:
            if hasattr(response, "result") and response.result:
                return self._extract_response_text(response.result)
        except Exception:
            pass

        # Final fallback: convert to string
        logger.error("Could not extract text from response using any method, using string conversion")
        return str(response)

    def _extract_json_array(self, text: str) -> list:
        """
        Extract JSON array from text that may contain extra content.
        Handles cases where there's text before/after the JSON array.
        """
        text = text.strip()

        # Try direct JSON parse first
        try:
            parsed = json.loads(text)
            if isinstance(parsed, list):
                return parsed
        except json.JSONDecodeError:
            pass

        # Find the first '[' and try to extract the array
        first_bracket = text.find("[")
        if first_bracket == -1:
            raise Exception("No JSON array found in response")

        # Try to find the matching closing bracket
        bracket_count = 0
        last_bracket = -1
        for i in range(first_bracket, len(text)):
            if text[i] == "[":
                bracket_count += 1
            elif text[i] == "]":
                bracket_count -= 1
                if bracket_count == 0:
                    last_bracket = i
                    break

        if last_bracket == -1:
            raise Exception("Could not find matching closing bracket for JSON array")

        # Extract just the array part
        json_text = text[first_bracket : last_bracket + 1]

        try:
            parsed = json.loads(json_text)
            if isinstance(parsed, list):
                logger.debug(f"Successfully extracted JSON array (length: {len(parsed)})")
                return parsed
            else:
                raise Exception(f"Extracted JSON is not an array, got: {type(parsed)}")
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse extracted JSON: {e}")
            logger.error(f"Extracted JSON text (first 500 chars): {json_text[:500]}")
            raise Exception(f"Failed to parse JSON array: {str(e)}")

    def _create_temp_dir(self) -> Path:
        """Create a temporary directory for processing."""
        self.temp_dir = Path(tempfile.mkdtemp())
        return self.temp_dir

    def _cleanup_temp_files(self):
        """Remove all temporary files and directories."""
        for file_path in self.temp_files:
            try:
                if file_path.exists():
                    if file_path.is_file():
                        file_path.unlink()
                    elif file_path.is_dir():
                        import shutil

                        shutil.rmtree(file_path)
            except Exception as e:
                print(f"Warning: Could not delete {file_path}: {e}")

        if self.temp_dir and self.temp_dir.exists():
            try:
                import shutil

                shutil.rmtree(self.temp_dir)
            except Exception as e:
                print(f"Warning: Could not delete temp dir {self.temp_dir}: {e}")

    def _resolve_soffice_cmd(self) -> str:
        """Resolve LibreOffice executable robustly (soffice may not be on PATH)."""
        soffice_cmd = (
            shutil.which("soffice")
            or shutil.which("libreoffice")
            or shutil.which("soffice.bin")
        )

        if not soffice_cmd:
            candidates = [
                "/usr/bin/soffice",
                "/usr/local/bin/soffice",
                "/usr/lib/libreoffice/program/soffice",
                "/usr/lib64/libreoffice/program/soffice",
                "/opt/libreoffice/program/soffice",
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            ]
            for c in candidates:
                if os.path.exists(c) and os.access(c, os.X_OK):
                    soffice_cmd = c
                    break

        if not soffice_cmd:
            raise FileNotFoundError(
                "LibreOffice executable not found. "
                "Install LibreOffice in the runtime image or ensure 'soffice'/'libreoffice' is on PATH."
            )

        return soffice_cmd

    def pptx_to_images(self, pptx_path: Path) -> Tuple[List[Path], List[str]]:
        """
        Convert PPTX slides to images and extract speaker notes.
        Returns tuple of (image_paths, speaker_notes_list).
        """
        try:
            logger.info(f"Loading PowerPoint file: {pptx_path}")
            prs = Presentation(pptx_path)
            total_slides = len(prs.slides)
            logger.info(f"Found {total_slides} slides in presentation")

            image_paths = []
            speaker_notes_list = []

            logger.info("Rendering PPTX to PNG using LibreOffice (soffice) ...")
            out_dir = self.temp_dir / "rendered_png"
            out_dir.mkdir(parents=True, exist_ok=True)
            self.temp_files.append(out_dir)

            soffice_cmd = self._resolve_soffice_cmd()

            # Try direct PPTX -> PNG export first
            subprocess.run(
                [
                    soffice_cmd,
                    "--headless",
                    "--nologo",
                    "--nolockcheck",
                    "--nodefault",
                    "--norestore",
                    "--convert-to",
                    "png:impress_png_Export",
                    "--outdir",
                    str(out_dir),
                    str(pptx_path),
                ],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                check=True,
            )

            png_files = list(out_dir.glob("*.png"))

            def _slide_sort_key(p: Path):
                m = re.search(r"(\d+)(?=\D*$)", p.stem)
                return int(m.group(1)) if m else 10**9

            png_files = sorted(png_files, key=_slide_sort_key)

            # ✅ FIX: If LibreOffice produced only 1 image (or fewer than slides), fallback to PDF -> PNG per page
            if len(png_files) < total_slides:
                logger.warning(
                    f"Direct PNG export produced {len(png_files)} image(s) for {total_slides} slides. "
                    "Falling back to PPTX -> PDF -> PNG per page."
                )

                pdftoppm_cmd = shutil.which("pdftoppm")
                if not pdftoppm_cmd:
                    raise FileNotFoundError(
                        "pdftoppm not found. Install 'poppler-utils' in the runtime image to enable PDF-to-PNG conversion."
                    )

                pdf_dir = self.temp_dir / "rendered_pdf"
                pdf_dir.mkdir(parents=True, exist_ok=True)
                self.temp_files.append(pdf_dir)

                pdf_out = pdf_dir / f"{pptx_path.stem}.pdf"
                self.temp_files.append(pdf_out)

                # PPTX -> PDF
                subprocess.run(
                    [
                        soffice_cmd,
                        "--headless",
                        "--nologo",
                        "--nolockcheck",
                        "--nodefault",
                        "--norestore",
                        "--convert-to",
                        "pdf:impress_pdf_Export",
                        "--outdir",
                        str(pdf_dir),
                        str(pptx_path),
                    ],
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    check=True,
                )

                if not pdf_out.exists():
                    pdf_candidates = list(pdf_dir.glob("*.pdf"))
                    if not pdf_candidates:
                        raise Exception("LibreOffice did not produce a PDF file for fallback rendering")
                    pdf_out = pdf_candidates[0]
                    self.temp_files.append(pdf_out)

                # PDF -> PNG per page
                out_prefix = str(out_dir / "slide")
                subprocess.run(
                    [
                        pdftoppm_cmd,
                        "-png",
                        "-r",
                        "200",
                        str(pdf_out),
                        out_prefix,
                    ],
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    check=True,
                )

                png_files = sorted(list(out_dir.glob("slide-*.png")), key=_slide_sort_key)

            if not png_files:
                raise Exception("No PNG files produced by either direct export or PDF fallback")

            if len(png_files) != total_slides:
                logger.warning(
                    f"Rendered PNG count ({len(png_files)}) does not match slide count ({total_slides})"
                )

            for i, slide in enumerate(prs.slides, start=1):
                logger.info(f"Processing slide {i}/{total_slides}...")

                notes = self._extract_speaker_notes(slide)
                if notes:
                    logger.debug(f"Slide {i} has speaker notes: {len(notes)} characters")
                else:
                    logger.debug(f"Slide {i} has no speaker notes")

                speaker_notes_list.append(notes)

            image_paths = [p for p in png_files]
            for p in image_paths:
                self.temp_files.append(p)

            logger.info(f"Successfully converted {len(image_paths)} slides to images")
            return image_paths, speaker_notes_list
        except Exception as e:
            logger.error(f"Failed to convert PPTX to images: {str(e)}")
            raise Exception(f"Failed to convert PPTX to images: {str(e)}")

    def _extract_slide_text(self, slide) -> str:
        """Extract text content from a slide."""
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_parts.append(text)
        return "\n".join(text_parts)

    def _extract_speaker_notes(self, slide) -> str:
        """Extract speaker notes from a slide."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    return notes_text if notes_text else ""
            return ""
        except Exception:
            # If notes slide doesn't exist or can't be accessed, return empty string
            return ""

    def _create_slide_image(self, text: str, output_path: Path) -> Image.Image:
        """
        Create a simple image representation of slide text.
        For MVP, this creates a text-based image.
        In production, you'd use libraries like python-pptx-render or convert via LibreOffice.
        """
        from PIL import Image, ImageDraw, ImageFont

        # Standard slide dimensions (16:9 aspect ratio)
        slide_width = 1920
        slide_height = 1080
        padding = 60

        # Create slide background
        img = Image.new("RGB", (slide_width, slide_height), color="#FFFFFF")
        draw = ImageDraw.Draw(img)

        # Try to load a better font, fallback to default
        try:
            # Try common system fonts
            font_large = ImageFont.truetype("arial.ttf", 48)
            font_normal = ImageFont.truetype("arial.ttf", 32)
        except Exception:
            try:
                font_large = ImageFont.truetype(
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48
                )
                font_normal = ImageFont.truetype(
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32
                )
            except Exception:
                font_large = ImageFont.load_default()
                font_normal = ImageFont.load_default()

        # Split text into lines and process
        lines = [line.strip() for line in text.split("\n") if line.strip()]
        if not lines:
            lines = ["[Empty slide]"]

        # Draw title (first line, larger)
        y = padding
        if lines:
            title = lines[0][:80]  # Limit title length
            draw.text((padding, y), title, fill="#1a1a1a", font=font_large)
            y += 80

        # Draw content (remaining lines, normal size)
        for line in lines[1:25]:  # Limit to 25 content lines
            if y + 50 > slide_height - padding:
                break
            # Handle long lines by wrapping
            if len(line) > 100:
                line = line[:97] + "..."
            draw.text((padding + 40, y), line, fill="#333333", font=font_normal)
            y += 45

        img.save(output_path, quality=95)
        return img

    def _get_style_instructions(self, narration_style: str) -> str:
        """Get style-specific instructions based on narration style."""
        styles = {
            "Human-like": """**Narration Style: Human-like**
Write as if you're speaking naturally to an audience:
- Use conversational transitions: "Now, let's look at...", "Moving forward...", "This brings us to..."
- Add natural connectives between slides: "Building on that idea...", "Following up on this point..."
- Reference specific points: "As you can see here...", "Notice how...", "This slide shows us..."
- Explain rather than just repeat: Don't read bullet points verbatim, explain their meaning
- Use natural pauses indicated by paragraph breaks for longer content
- Make it sound like you're genuinely reading and explaining the slides""",
            "Formal": """**Narration Style: Formal**
Write in a formal, structured manner:
- Use formal transitions: "We will now examine...", "Proceeding to...", "Let us consider..."
- Maintain professional language throughout
- Avoid contractions and casual expressions
- Use structured transitions: "Furthermore...", "Moreover...", "In addition..."
- Present information systematically and authoritatively
- Keep transitions professional and clear""",
            "Concise": """**Narration Style: Concise**
Write brief, to-the-point narration:
- Get straight to the point, avoid unnecessary words
- Use direct transitions: "Next...", "Moving to...", "Now..."
- Focus on key points only
- Keep sentences short and clear
- Minimize filler words and phrases
- Be efficient with transitions between slides""",
            "Storytelling": """**Narration Style: Storytelling**
Write as a narrative that tells a story:
- Create a narrative arc across slides
- Use storytelling transitions: "Our journey begins with...", "As the story unfolds...", "This brings us to a pivotal moment..."
- Build connections between slides like chapters in a story
- Use descriptive language to paint a picture
- Create anticipation: "What we'll discover next...", "This sets the stage for..."
- Make the presentation feel like a cohesive narrative""",
            "Conversational": """**Narration Style: Conversational**
Write in a friendly, approachable conversational style:
- Use casual, friendly transitions
- Include rhetorical questions
- Use everyday language and relatable examples
- Create a dialogue feel
- Make transitions feel like natural conversation flow""",
            "Professional": """**Narration Style: Professional**
Write in a polished, business-appropriate style:
- Use professional transitions: "We'll now explore...", "Turning our attention to...", "Let's examine..."
- Maintain a balanced, confident tone
- Use clear, structured language
- Include appropriate business terminology
- Create smooth, logical transitions
- Present information with authority and clarity""",
        }
        return styles.get(narration_style, styles["Human-like"])

    def _get_length_instructions(self, dynamic_length: bool) -> str:
        """Get length-specific instructions."""
        if dynamic_length:
            return """**Dynamic Length**: Adjust the narration length based on slide content complexity.
- You are responsible for determining the slide’s complexity before writing the narration
- Simple slides (low complexity): 50–100 words, concise and clear
- Medium complexity slides: 100–150 words, with added explanation
- Complex slides (high complexity): 150–200 words, structured into multiple paragraphs if needed
- Use the Content Complexity indicator to guide narration length
- For longer narrations, only when necessary use 200-400 words), split into 2–3 paragraphs using "\\n\\n" (double newline)
- Never exceed 400 words under any circumstances
"""

        else:
            return """**Fixed Length**: Keep narration consistent across slides:
- Aim for 100-150 words per slide
- Maintain similar length for all slides regardless of content complexity
- Break into paragraphs only if exceeding 200 words
- Use "\\n\\n" (double newline) for paragraph breaks when needed"""

    def process_slide_with_gemini(
        self,
        image_path: Path,
        slide_number: int,
        tone: str,
        audience_level: str,
        max_words: int,
    ) -> str:
        """
        Process a single slide image with Gemini to get rewritten content only.
        Returns rewritten content string.
        """
        try:
            logger.info(f"Processing slide {slide_number} with Gemini...")

            # Read image as PIL Image object (required by Gemini API)
            image = Image.open(image_path)

            # Build prompt - only ask for rewritten content, not speaker notes
            prompt = f"""You are an expert presentation script writer. Analyze this slide image and create a clear, engaging narration script that explains the content on the slide.

- Make it structured, clear, and concise (max {max_words} words)
- Explain the slide content meaningfully.
- Focus on explaining diagrams and tables. Only include and describe those that are directly useful and relevant to the main information. Ignore images, icons, or visuals used solely for aesthetics
- Maintain the key information and meaning
- Tone: {tone}
- Audience: {audience_level}
- Slide number: {slide_number}

- Return your response as a JSON object with exactly this key:
{{
    "rewritten_content": "narration script explaining slide content here"
}}

- CRITICAL:
- Only return valid JSON, no markdown formatting or additional text
- The "rewritten_content" value must be plain text only – NO markdown formatting, NO markdown syntax (no **, *, _, #, [], etc.), NO special formatting characters. Use only plain text"""


            # Call Gemini with image (PIL Image object)
            logger.info(f"Calling Gemini API for slide {slide_number}...")
            response = self.model.generate_content([prompt, image])

            # Parse JSON response - safely extract text
            response_text = self._extract_response_text(response).strip()
            logger.debug(f"Raw Gemini response for slide {slide_number}: {response_text[:200]}...")

            # Remove markdown code blocks if present
            if response_text.startswith("```"):
                response_text = response_text.split("```")[1]
                if response_text.startswith("json"):
                    response_text = response_text[4:]
                response_text = response_text.strip()

            result = self._safe_json_loads(response_text)

            # Extract rewritten_content and ensure it's a string
            rewritten_content = result.get("rewritten_content", "")

            # Handle case where Gemini returns a list instead of string
            if isinstance(rewritten_content, list):
                logger.warning(
                    f"Slide {slide_number}: Gemini returned list instead of string, converting..."
                )
                rewritten_content = "\n".join(str(item) for item in rewritten_content)
            elif not isinstance(rewritten_content, str):
                logger.warning(
                    f"Slide {slide_number}: Unexpected type {type(rewritten_content)}, converting to string..."
                )
                rewritten_content = str(rewritten_content)

            logger.info(f"Successfully processed slide {slide_number}")
            return rewritten_content

        except json.JSONDecodeError as e:
            logger.error(f"JSON decode error for slide {slide_number}: {str(e)}")
            logger.error(
                f"Response text: {response_text[:500] if 'response_text' in locals() else 'N/A'}"
            )
            raise Exception(f"Failed to parse Gemini response as JSON: {str(e)}")
        except Exception as e:
            logger.error(f"Error processing slide {slide_number}: {str(e)}")
            raise Exception(f"Gemini API error: {str(e)}")

    # ---------------------------
    # NARRATION FIX STARTS HERE
    # ---------------------------

    def _extract_first_json_object(self, text: str) -> str:
        """
        Extract the first top-level JSON object from a string.
        This avoids failures when the model returns extra text around JSON.
        """
        text = text.strip()
        start = text.find("{")
        if start == -1:
            raise Exception("No JSON object found in response")

        depth = 0
        in_string = False
        escape = False

        for i in range(start, len(text)):
            ch = text[i]

            if in_string:
                if escape:
                    escape = False
                elif ch == "\\":
                    escape = True
                elif ch == '"':
                    in_string = False
                continue
            else:
                if ch == '"':
                    in_string = True
                    continue
                if ch == "{":
                    depth += 1
                elif ch == "}":
                    depth -= 1
                    if depth == 0:
                        return text[start : i + 1]

        raise Exception("Could not find matching closing brace for JSON object")

    def _clean_json_control_chars(self, s: str) -> str:
        """
        Remove unescaped control characters that break json.loads.
        This targets characters in U+0000..U+001F (except valid JSON escapes).
        The most common culprit from LLM output is a literal newline or tab inside a JSON string.
        We do NOT want to change intended paragraph breaks. So:
        - We do not blindly remove '\\n' (two characters). That is valid JSON.
        - We remove actual control characters: '\n', '\r', '\t', and other <0x20 chars.
        """
        if not s:
            return s

        out = []
        for ch in s:
            code = ord(ch)
            if code < 32:
                # drop actual control chars; json requires them to be escaped (\\n, \\t, etc.)
                continue
            out.append(ch)
        return "".join(out)

    def _safe_json_loads(self, response_text: str) -> Dict:
        """
        Robust JSON parse:
        1) Strip code fences if present
        2) Try direct json.loads
        3) Extract first JSON object and try again
        4) Remove control chars and try again (both full text and extracted object)
        """
        text = (response_text or "").strip()

        # Remove markdown code blocks if present
        if text.startswith("```"):
            parts = text.split("```")
            if len(parts) >= 2:
                text = parts[1]
                if text.startswith("json"):
                    text = text[4:]
                text = text.strip()

        # Attempt direct parse
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass

        # Try extracting first JSON object
        try:
            candidate = self._extract_first_json_object(text)
            try:
                return json.loads(candidate)
            except json.JSONDecodeError:
                # try cleaned candidate
                cleaned_candidate = self._clean_json_control_chars(candidate)
                return json.loads(cleaned_candidate)
        except Exception:
            pass

        # As a last resort, clean the full text and parse again
        cleaned = self._clean_json_control_chars(text)
        return json.loads(cleaned)

    def _compute_complexity_label(self, slide_content: str, speaker_notes: str) -> str:
        """Compute Low/Medium/High complexity based on current slide only."""
        content_length = len(slide_content.split()) if slide_content else 0
        notes_length = len(speaker_notes.split()) if speaker_notes else 0
        total = content_length + notes_length
        return "High" if total > 150 else "Medium" if total > 50 else "Low"

    def _length_target_hint(self, dynamic_length: bool, complexity: str) -> str:
        """Provide a per-slide word target hint without changing global behavior."""
        if not dynamic_length:
            return "Aim for 100-150 words."
        if complexity == "Low":
            return "Aim for 50-100 words."
        if complexity == "Medium":
            return "Aim for 100-200 words."
        return "Aim for 200-400 words, split into 2-3 paragraphs using \\n\\n if needed."

    def _generate_single_slide_narration(
        self,
        slide_index: int,
        total_slides: int,
        slide_content: str,
        speaker_notes: str,
        prev_narrations: List[str],
        tone: str,
        narration_style: str,
        dynamic_length: bool,
    ) -> str:
        """
        Generate narration for ONE slide only, using only the past 3 narrations as context.
        Returns narration as plain text (single string).
        """
        slide_number = slide_index + 1

        style_instructions = self._get_style_instructions(narration_style)
        length_instructions = self._get_length_instructions(dynamic_length)

        complexity = self._compute_complexity_label(slide_content, speaker_notes)
        per_slide_hint = self._length_target_hint(dynamic_length, complexity)

        # Only last 5 narrations as requested
        prev_narrations = (prev_narrations or [])[-5:]

        if prev_narrations:
            prev_block = "\n\n".join(
                [f"- Prior narration {len(prev_narrations) - i}: {n}" for i, n in enumerate(prev_narrations)]
            )
        else:
            prev_block = "[None]"

        needs_opening_connection = slide_number != 1
        needs_closing_transition = slide_number != total_slides

        prompt = f"""You are a professional presenter creating a {narration_style.lower()} narration script.

{style_instructions}

{length_instructions}

Tone: Maintain a {tone} tone throughout.

IMPORTANT CONTEXT RULES:
- You may ONLY use the past 5 narrations provided below as cross-slide context.
- Do NOT invent, reference, or imply any other slides beyond what is provided.
- Generally avoid repeating the same phrases, sentence structures, or opening flows from previous narrations.
- Reuse wording from past narrations only when it is necessary for clarity or continuity, and do not overuse it.

Past narrations (most recent last):
{prev_block}

Current slide to narrate:
- Slide number: {slide_number} of {total_slides}
- Rewritten Content:
{slide_content}
- Speaker Notes:
{speaker_notes}

Structure requirements for THIS slide:
- Start in a way that fits the context, but vary the opening so it does not feel repetitive.
- Use a transition from previous narrations only when it adds value; avoid forced connectors.
- Explain the slide content meaningfully (do not read or restate bullets verbatim).
- Incorporate relevant speaker notes naturally, only when they add value and context.
- {"End with a transition to the next slide only if clearly suggested by the speaker notes." if needs_closing_transition else "Do NOT add a transition to a next slide; close the narration naturally."}

Return your response as a JSON object with exactly this key:
{{
  "narration": "plain text narration here"
}}

CRITICAL: 
- Only return valid JSON, no markdown, no extra text. 
- The narration must be plain text only 
- NO markdown formatting, NO markdown syntax (no **, *, _, #, [], (), etc.), NO code blocks. 
- If you need paragraph breaks, represent them using the two-character sequence \\n\\n (backslash-n-backslash-n) inside the JSON string. 
- Do NOT include literal newlines or literal tabs inside the JSON string value (they must be escaped as \\n and \\t).
"""

        logger.info(f"Calling Gemini API for narration (slide {slide_number}/{total_slides})...")
        response = self.model.generate_content(prompt)

        response_text = self._extract_response_text(response).strip()
        logger.debug(
            f"Raw single-slide narration response (slide {slide_number}): {response_text[:300]}..."
        )

        parsed = self._safe_json_loads(response_text)

        narration = parsed.get("narration", "")

        if isinstance(narration, list):
            logger.warning(f"Slide {slide_number}: narration returned as list; joining.")
            narration = "\n".join(str(x) for x in narration)
        elif not isinstance(narration, str):
            narration = str(narration)

        narration = narration.strip()
        narration = narration.replace("\\n\\n", "\n\n").replace("\\n", "\n").replace("\\t", "\t")

        return narration

    def generate_narration(
        self,
        slide_data: List[Dict[str, str]],
        tone: str,
        narration_style: str = "Human-like",
        dynamic_length: bool = True,
    ) -> List[str]:
        """
        Generate narration with configurable style and length options.

        Generates narration PER SLIDE (separate Gemini calls),
        uses ONLY the past 5 narrations as cross-slide context.
        """
        try:
            logger.info(f"Generating narration for {len(slide_data)} slides...")
            logger.info(f"Narration style: {narration_style}, Dynamic length: {dynamic_length}")

            processed_narrations: List[str] = []
            total = len(slide_data)

            for i, slide in enumerate(slide_data):
                slide_content = slide.get("rewritten_content", "") or ""
                speaker_notes = slide.get("speaker_notes", "") or ""

                prev_3 = processed_narrations[-3:]  # ONLY past 3 narrations as context

                narration = self._generate_single_slide_narration(
                    slide_index=i,
                    total_slides=total,
                    slide_content=slide_content,
                    speaker_notes=speaker_notes,
                    prev_narrations=prev_3,
                    tone=tone,
                    narration_style=narration_style,
                    dynamic_length=dynamic_length,
                )

                word_count = len(narration.split()) if narration else 0
                logger.info(f"Slide {i+1} narration ({word_count} words): {narration[:100]}...")
                processed_narrations.append(narration)

            expected_count = len(slide_data)
            if len(processed_narrations) != expected_count:
                logger.error(
                    f"CRITICAL: Processed narrations count ({len(processed_narrations)}) "
                    f"doesn't match slide count ({expected_count})"
                )
                while len(processed_narrations) < expected_count:
                    processed_narrations.append("")
                processed_narrations = processed_narrations[:expected_count]

            total_words = sum(len(n.split()) for n in processed_narrations)
            avg_words = total_words / len(processed_narrations) if processed_narrations else 0
            logger.info(f"Successfully generated {len(processed_narrations)} narration paragraphs (one per slide)")
            logger.info(f"Total words: {total_words}, Average per slide: {avg_words:.1f} words")

            return processed_narrations

        except Exception as e:
            logger.error(f"Failed to generate narration: {str(e)}")
            raise Exception(f"Failed to generate narration: {str(e)}")

    # ---------------------------
    # NARRATION FIX ENDS HERE
    # ---------------------------

    def process_pptx(
        self,
        pptx_path: Path,
        tone: str,
        audience_level: str,
        max_words: int,
        notes_length: str,
        narration_style: str = "Human-like",
        dynamic_length: bool = True,
    ) -> Dict:
        """
        Main processing pipeline: convert PPTX, process slides, generate narration.
        """
        try:
            logger.info("=" * 60)
            logger.info("Starting PPTX processing pipeline")
            logger.info(f"File: {pptx_path}")
            logger.info(f"Tone: {tone}, Audience: {audience_level}, Max words: {max_words}")
            logger.info("=" * 60)

            self._create_temp_dir()
            self.temp_files.append(pptx_path)
            logger.info(f"Created temp directory: {self.temp_dir}")

            logger.info("\n[STEP 1] Converting PPTX to images and extracting speaker notes...")
            image_paths, speaker_notes_list = self.pptx_to_images(pptx_path)

            if not image_paths:
                raise Exception("No slides found in presentation")

            logger.info(f"✓ Converted {len(image_paths)} slides to images\n")

            logger.info(f"[STEP 2] Processing {len(image_paths)} slides with Gemini...")
            slide_results = []
            for i, img_path in enumerate(image_paths, start=1):
                logger.info(f"\n--- Processing Slide {i}/{len(image_paths)} ---")
                rewritten_content = self.process_slide_with_gemini(
                    img_path, i, tone, audience_level, max_words
                )

                if isinstance(rewritten_content, list):
                    logger.warning(f"Slide {i}: Converting list to string")
                    rewritten_content = "\n".join(str(item) for item in rewritten_content)
                elif not isinstance(rewritten_content, str):
                    logger.warning(f"Slide {i}: Converting {type(rewritten_content)} to string")
                    rewritten_content = str(rewritten_content)

                speaker_notes = speaker_notes_list[i - 1] if i - 1 < len(speaker_notes_list) else ""

                if not isinstance(speaker_notes, str):
                    speaker_notes = str(speaker_notes)

                slide_results.append(
                    {
                        "slide_number": i,
                        "rewritten_content": rewritten_content,
                        "speaker_notes": speaker_notes,
                    }
                )

                logger.info(f"✓ Slide {i} processed successfully")
                logger.info(f"  Rewritten content length: {len(rewritten_content)} chars")
                logger.info(f"  Speaker notes length: {len(speaker_notes)} chars")

            logger.info(f"\n✓ All {len(slide_results)} slides processed\n")

            logger.info("[STEP 3] Generating flowing narration...")
            narration_paragraphs = self.generate_narration(
                slide_results,
                tone,
                narration_style=narration_style,
                dynamic_length=dynamic_length,
            )

            if len(narration_paragraphs) != len(slide_results):
                logger.error(
                    f"CRITICAL ERROR: Narration count ({len(narration_paragraphs)}) "
                    f"doesn't match slide count ({len(slide_results)})"
                )
                while len(narration_paragraphs) < len(slide_results):
                    narration_paragraphs.append("")
                narration_paragraphs = narration_paragraphs[: len(slide_results)]
                logger.warning(f"Adjusted narration array to match {len(slide_results)} slides")

            logger.info(
                f"✓ Generated {len(narration_paragraphs)} narration paragraphs (validated: one per slide)\n"
            )

            logger.info("[STEP 4] Combining results with strict slide-to-narration mapping...")
            final_results = []
            for i, slide_data in enumerate(slide_results):
                rewritten_content = slide_data["rewritten_content"]
                speaker_notes = slide_data["speaker_notes"]

                narration = narration_paragraphs[i] if i < len(narration_paragraphs) else ""

                if isinstance(rewritten_content, list):
                    rewritten_content = "\n".join(str(item) for item in rewritten_content)
                elif not isinstance(rewritten_content, str):
                    rewritten_content = str(rewritten_content)

                if not isinstance(speaker_notes, str):
                    speaker_notes = str(speaker_notes)

                if not isinstance(narration, str):
                    narration = str(narration)

                final_results.append(
                    {
                        "slide_number": slide_data["slide_number"],
                        "rewritten_content": rewritten_content,
                        "speaker_notes": speaker_notes,
                        "narration_paragraph": narration,
                    }
                )

            logger.info("=" * 60)
            logger.info("Processing complete!")
            logger.info(f"Total slides processed: {len(final_results)}")
            logger.info("=" * 60)

            return {
                "success": True,
                "total_slides": len(final_results),
                "slides": final_results,
                "overall_summary": None,
            }
        except Exception as e:
            logger.error("=" * 60)
            logger.error(f"ERROR in processing pipeline: {str(e)}")
            logger.error("=" * 60)
            import traceback

            logger.error(traceback.format_exc())
            return {"success": False, "total_slides": 0, "slides": [], "error": str(e)}
        finally:
            logger.info("Cleaning up temporary files...")
            self._cleanup_temp_files()
            logger.info("Cleanup complete")
